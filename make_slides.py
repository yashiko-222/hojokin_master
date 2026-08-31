"""
補助金マッチングツール プレゼン資料（.pptx）生成スクリプト
構成: タイトル / アプリの目的 / 現状分析・課題 / 導入効果 / アーキテクチャ図
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn


# ===== カラーパレット =====
NAVY = RGBColor(0x1F, 0x3A, 0x5F)
BLUE = RGBColor(0x15, 0x65, 0xC0)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
LIGHT_ORANGE = RGBColor(0xFD, 0xEB, 0xD0)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_RED = RGBColor(0xFB, 0xE3, 0xE0)
GRAY = RGBColor(0x60, 0x60, 0x60)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x22, 0x22, 0x22)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


def add_textbox(slide, left, top, width, height, text, size=18,
                bold=False, color=DARK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font="Meiryo", space_after=6):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def add_box(slide, left, top, width, height, text, fill, line_color,
            text_color=WHITE, size=14, bold=True, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    box = slide.shapes.add_shape(shape, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = fill
    box.line.color.rgb = line_color
    box.line.width = Pt(1.25)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold if i == 0 else False
        run.font.color.rgb = text_color
        run.font.name = "Meiryo"
    return box


def add_arrow(slide, x1, y1, x2, y2, color=NAVY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2.25)
    ln = conn.line._get_or_add_ln()
    tail = ln.makeelement(qn("a:tailEnd"),
                          {"type": "triangle", "w": "med", "len": "med"})
    ln.append(tail)
    return conn


def slide_header(slide, title, subtitle=None):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(1.1))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    add_textbox(slide, Inches(0.5), Inches(0.12), Inches(12.3), Inches(0.6),
                title, size=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(0.72), Inches(12.3), Inches(0.35),
                    subtitle, size=13, color=LIGHT_BLUE)


def card(slide, x, y, w, h, title, body, fill, ln, title_size=16, body_size=13):
    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = fill
    c.line.color.rgb = ln
    c.line.width = Pt(1.5)
    add_textbox(slide, Emu(int(x) + int(Inches(0.25))),
                Emu(int(y) + int(Inches(0.2))),
                Emu(int(w) - int(Inches(0.5))), Inches(0.5),
                title, size=title_size, bold=True, color=ln)
    add_textbox(slide, Emu(int(x) + int(Inches(0.25))),
                Emu(int(y) + int(Inches(0.8))),
                Emu(int(w) - int(Inches(0.5))), Emu(int(h) - int(Inches(1.0))),
                body, size=body_size, color=DARK)


# ============================================================
# スライド1: タイトル
# ============================================================
s = prs.slides.add_slide(BLANK)
bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(4.5), SW, Inches(0.12))
accent.fill.solid(); accent.fill.fore_color.rgb = RGBColor(0x4F, 0x9E, 0xE3)
accent.line.fill.background()
add_textbox(s, Inches(0.9), Inches(2.4), Inches(11.5), Inches(1.2),
            "補助金マッチングツール", size=48, bold=True, color=WHITE)
add_textbox(s, Inches(0.95), Inches(3.6), Inches(11.5), Inches(0.8),
            "企業HPを解析し、最適な補助金を自動提案するアプリ",
            size=20, color=LIGHT_BLUE)
add_textbox(s, Inches(0.95), Inches(5.0), Inches(11.5), Inches(1.2),
            "目的 ／ 現状分析・課題 ／ 導入効果 ／ アーキテクチャ",
            size=15, color=WHITE)

# ============================================================
# スライド2: アプリの目的
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "アプリの目的", "自社に合う補助金を、誰でも・すぐに・見つけられるようにする")

add_textbox(s, Inches(0.7), Inches(1.4), Inches(12), Inches(0.8),
            "企業HPのURLまたは企業名を入力するだけで、事業内容を自動で分析し、"
            "最適な補助金候補を提示することを目的とします。",
            size=16, color=DARK)

goals = [
    ("探す手間をなくす", "膨大な補助金制度から\n自社に合うものを\n人手で探す負担を削減", LIGHT_BLUE, BLUE),
    ("専門知識が不要", "事業内容を自動で解析し\n候補を提示するため\n知識がなくても使える", LIGHT_GREEN, GREEN),
    ("機会損失を防ぐ", "見落としがちな補助金を\n提示し、活用できる\n制度の取りこぼしを防止", LIGHT_ORANGE, ORANGE),
]
w = Inches(3.9); h = Inches(2.6); y = Inches(2.7); x0 = Inches(0.6); gap = Inches(0.35)
for i, (t, b, fill, ln) in enumerate(goals):
    x = Emu(int(x0) + i * (int(w) + int(gap)))
    card(s, x, y, w, h, t, b, fill, ln)

add_textbox(s, Inches(0.7), Inches(5.6), Inches(12), Inches(1.2),
            "対象ユーザー：補助金の活用を検討する中小企業・小規模事業者、"
            "および支援を行う担当者\n"
            "データ元：デジタル庁「Jグランツ」を中心とした公的補助金情報",
            size=13, color=GRAY)

# ============================================================
# スライド3: 現状分析・課題
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "現状分析・課題", "補助金は多数あるが「自社に合うものを見つける」ことが難しい")

add_textbox(s, Inches(0.7), Inches(1.35), Inches(12), Inches(0.5),
            "■ 現状", size=16, bold=True, color=NAVY)
add_textbox(s, Inches(0.9), Inches(1.85), Inches(11.8), Inches(1.0),
            "・国や自治体が多数の補助金・助成金を公開しているが、制度が多く複雑\n"
            "・名称変更や公募時期の更新が頻繁で、最新情報の把握が難しい",
            size=14, color=DARK)

add_textbox(s, Inches(0.7), Inches(3.0), Inches(12), Inches(0.5),
            "■ 課題", size=16, bold=True, color=RED)

issues = [
    ("探索コスト", "自社に合う補助金を\n探すのに時間と\n手間がかかる"),
    ("専門性の壁", "対象要件や分野の\n判断に専門知識が\n必要になる"),
    ("情報の陳腐化", "URLや制度名が変わり\n古い情報にたどり\n着いてしまう"),
    ("機会損失", "使えるはずの補助金を\n知らずに\n見逃してしまう"),
]
w = Inches(2.9); h = Inches(2.0); y = Inches(3.55); x0 = Inches(0.6); gap = Inches(0.2)
for i, (t, b) in enumerate(issues):
    x = Emu(int(x0) + i * (int(w) + int(gap)))
    card(s, x, y, w, h, t, b, LIGHT_RED, RED, title_size=15, body_size=12)

add_textbox(s, Inches(0.7), Inches(5.85), Inches(12), Inches(1.0),
            "→ 「事業内容の分析」と「補助金情報とのマッチング」を自動化することで、"
            "これらの課題を解決する。",
            size=15, bold=True, color=NAVY)

# ============================================================
# スライド4: 導入効果
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "導入効果", "探索の自動化により、時間・精度・活用機会を改善")

effects = [
    ("時間の短縮", "URL/企業名を入れるだけで\n候補が一覧化され、\n探索時間を大幅に短縮", LIGHT_BLUE, BLUE),
    ("精度の向上", "キーワードに加え事業概要・\n業種も分析し、関連度順に\n根拠（理由）付きで提示", LIGHT_GREEN, GREEN),
    ("取りこぼし防止", "複数候補を提示するため\n見落としを防ぎ、\n活用機会を最大化", LIGHT_ORANGE, ORANGE),
]
w = Inches(3.9); h = Inches(2.4); y = Inches(1.55); x0 = Inches(0.6); gap = Inches(0.35)
for i, (t, b, fill, ln) in enumerate(effects):
    x = Emu(int(x0) + i * (int(w) + int(gap)))
    card(s, x, y, w, h, t, b, fill, ln)

# Before / After
add_textbox(s, Inches(0.7), Inches(4.3), Inches(12), Inches(0.5),
            "Before → After", size=16, bold=True, color=NAVY)
ba = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        Inches(0.6), Inches(4.85), Inches(12.1), Inches(2.0))
ba.fill.solid(); ba.fill.fore_color.rgb = LIGHT_GRAY
ba.line.color.rgb = GRAY
add_textbox(s, Inches(0.9), Inches(5.0), Inches(11.6), Inches(1.7),
            "従来：担当者が複数サイトを横断し、制度を1つずつ調べて該当性を判断（時間・知識が必要）\n"
            "本アプリ：企業情報を入力 → 自動で事業を分析 → 最適な補助金候補を関連度順に提示\n"
            "効果：探索の属人化を解消し、誰でも短時間で「使える補助金の当たり」をつけられる",
            size=13, color=DARK)

# ============================================================
# スライド5: アーキテクチャ図
# ============================================================
s = prs.slides.add_slide(BLANK)
slide_header(s, "システムアーキテクチャ",
             "フロントエンド（Streamlit）と バックエンド（FastAPI）を HTTP API で連携")

row_y = Inches(2.0); row_h = Inches(1.0)
add_box(s, Inches(0.6), row_y, Inches(2.4), row_h,
        "フロントエンド\nStreamlit (app.py)\nブラウザUI", BLUE, NAVY, size=13)
add_box(s, Inches(3.9), row_y, Inches(2.6), row_h,
        "バックエンド\nFastAPI (backend/main.py)\n/api/search・/api/analyze",
        GREEN, RGBColor(0x1B, 0x5E, 0x20), size=12)

mod_x = Inches(7.4); mod_w = Inches(2.5); mod_h = Inches(0.62)
modules = [
    ("search.py", "企業名→公式HP特定"),
    ("crawler.py", "クロール・キーワード/概要/業種"),
    ("matcher.py", "補助金マッチング(関連度)"),
    ("jgrants.py", "補助金データ"),
]
mod_top0 = Inches(1.55); gap = Inches(0.12)
for i, (name, desc) in enumerate(modules):
    top = Emu(int(mod_top0) + i * (int(mod_h) + int(gap)))
    add_box(s, mod_x, top, mod_w, mod_h, f"{name} … {desc}",
            LIGHT_GREEN, GREEN, text_color=DARK, size=10.5)

ext_x = Inches(10.5); ext_w = Inches(2.3); ext_h = Inches(0.75)
externals = ["Bing 検索", "Wikipedia", "Jグランツ / 公式サイト"]
ext_top0 = Inches(1.55); ext_gap = Inches(0.2)
for i, name in enumerate(externals):
    top = Emu(int(ext_top0) + i * (int(ext_h) + int(ext_gap)))
    add_box(s, ext_x, top, ext_w, ext_h, name, LIGHT_ORANGE, ORANGE,
            text_color=DARK, size=11)

add_textbox(s, mod_x, Inches(1.2), mod_w, Inches(0.3),
            "処理モジュール (modules/)", size=11, bold=True, color=GREEN)
add_textbox(s, ext_x, Inches(1.2), ext_w, Inches(0.3),
            "外部データソース", size=11, bold=True, color=ORANGE)

add_arrow(s, Inches(3.0), Inches(2.5), Inches(3.9), Inches(2.5))
add_textbox(s, Inches(2.85), Inches(2.05), Inches(1.2), Inches(0.4),
            "HTTP\n(JSON)", size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_arrow(s, Inches(6.5), Inches(2.5), Inches(7.4), Inches(2.0))
add_arrow(s, Inches(9.9), Inches(1.9), Inches(10.5), Inches(1.9))
add_arrow(s, Inches(9.9), Inches(3.0), Inches(10.5), Inches(3.3))

flow = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                          Inches(0.6), Inches(4.6), Inches(12.1), Inches(2.3))
flow.fill.solid(); flow.fill.fore_color.rgb = LIGHT_GRAY
flow.line.color.rgb = GRAY
add_textbox(s, Inches(0.9), Inches(4.75), Inches(11.6), Inches(0.4),
            "処理の流れ", size=15, bold=True, color=NAVY)
add_textbox(s, Inches(0.9), Inches(5.2), Inches(11.6), Inches(1.6),
            "1. 入力判定：URLならそのまま／企業名なら search.py が Bing・Wikipedia で公式HP候補を検索\n"
            "2. クロール：crawler.py が HP を取得し、キーワード・事業概要・推定業種を抽出\n"
            "3. マッチング：matcher.py が概要・業種を加味し、jgrants.py の補助金を関連度順にランキング\n"
            "4. 結果表示：Streamlit が概要・業種・キーワード・補助金候補（推薦理由付き）を表示",
            size=13, color=DARK)

prs.save("補助金マッチングツール_紹介資料.pptx")
print("saved")
